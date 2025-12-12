#!/usr/bin/env python3
"""
IRIS Bicycle Incidents Pipeline - Streamlit Application
Complete pipeline: Excel to Geocoding to Clustering to Maps to PowerPoint
Barcelona Mobility Analysis Tool
"""

import streamlit as st
import pandas as pd
import geopandas as gpd
import numpy as np
import os
import re
import time
import tempfile
from pathlib import Path
from collections import defaultdict
from datetime import datetime
from io import BytesIO

# Geospatial imports
import osmnx as ox
from shapely.geometry import Point
from sklearn.cluster import DBSCAN, KMeans

# Visualization
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import contextily as cx
from matplotlib.patches import Rectangle, Patch

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

# ============================================================================
# PAGE CONFIGURATION
# ============================================================================

st.set_page_config(
    page_title="IRIS Pipeline - Barcelona",
    page_icon="",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================================
# PROFESSIONAL WHITE THEME
# ============================================================================

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
    
    :root {
        --primary: #1a1a1a;
        --secondary: #4a4a4a;
        --accent: #0066cc;
        --border: #e0e0e0;
        --bg-white: #ffffff;
        --bg-light: #f8f9fa;
        --success: #28a745;
    }
    
    .stApp {
        background-color: var(--bg-white);
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
    }
    
    h1, h2, h3, h4, h5, h6 {
        font-family: 'Inter', sans-serif !important;
        color: var(--primary) !important;
        font-weight: 600 !important;
    }
    
    h1 {
        font-size: 1.75rem !important;
        border-bottom: 2px solid var(--border);
        padding-bottom: 0.75rem;
        margin-bottom: 1.5rem !important;
    }
    
    section[data-testid="stSidebar"] {
        background-color: var(--bg-light);
        border-right: 1px solid var(--border);
    }
    
    div[data-testid="metric-container"] {
        background: var(--bg-light);
        border: 1px solid var(--border);
        border-radius: 6px;
        padding: 1rem;
    }
    
    div[data-testid="stFileUploader"] {
        background: var(--bg-light);
        border: 1px dashed var(--border);
        border-radius: 6px;
        padding: 1.5rem;
    }
    
    .stButton > button {
        background-color: var(--accent) !important;
        color: white !important;
        border: none !important;
        border-radius: 4px !important;
        font-weight: 500 !important;
    }
    
    .stDownloadButton > button {
        background-color: var(--success) !important;
        color: white !important;
        border: none !important;
        border-radius: 4px !important;
    }
    
    .stProgress > div > div {
        background-color: var(--accent) !important;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        background: var(--bg-light);
        border-radius: 4px;
        padding: 2px;
        border: 1px solid var(--border);
    }
    
    .stTabs [data-baseweb="tab"] {
        background: transparent !important;
        color: var(--secondary) !important;
        font-weight: 500 !important;
    }
    
    .stTabs [aria-selected="true"] {
        background: var(--bg-white) !important;
        color: var(--primary) !important;
    }
</style>
""", unsafe_allow_html=True)


# ============================================================================
# CONSTANTS (exact from notebook)
# ============================================================================

A3_WIDTH_IN = 16.54
A3_HEIGHT_IN = 11.69
MAP_FRAC = 0.75
BOX_FRAC = 0.25

BARCELONA_MTM_TILES = 'https://cdn-geo.bcn.cat/XYZ/MTM/{z}/{x}/{y}.png'

COUNT_COLORS = {
    '5+': '#d62728',
    '4': '#ff7f0e',
    '3': '#ffcc00',
    '2': '#2ca02c',
    '1': '#1f77b4'
}

COUNT_ORDER = ['5+', '4', '3', '2', '1']


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def expand_street_type(abbrev):
    abbrev_map = {
        'Pl': 'Placa', 'Plta': 'Placeta', 'Pg': 'Passeig',
        'Av': 'Avinguda', 'Rbla': 'Rambla', 'Pla': 'Pla'
    }
    return abbrev_map.get(str(abbrev).strip(), abbrev) if pd.notna(abbrev) else None


def clean_street_name(street_name):
    if pd.isna(street_name):
        return ""
    street = str(street_name).strip()
    street = re.sub(r'^(Av|Pg|Pl|Plta|Rbla|Pla)\s+', '', street)
    street = re.sub(r'(\w+)\s+(?:[A-Z]\.\s+)+([A-Z]\w+)', r'\1 \2', street)
    return re.sub(r'\s+', ' ', street).strip()


def create_complete_address(row, street_col='carrer', num_col='numero_inici', 
                             barri_col='barri', districte_col='districte'):
    street_type = None
    if 'tipus_via' in row and pd.notna(row.get('tipus_via')):
        street_type = expand_street_type(row['tipus_via'])
    
    if pd.notna(row.get(street_col)):
        street_name = clean_street_name(row[street_col])
        if not street_type and street_name.startswith(('de ', 'del ', 'dels ', "de l'")):
            street_type = 'Carrer'
        full_street = f"{street_type} {street_name}".strip() if street_type else street_name
    else:
        return ""
    
    if pd.notna(row.get(num_col)) and row[num_col] != 0:
        number = str(int(row[num_col]))
        if 'lletra_inici' in row and pd.notna(row.get('lletra_inici')):
            letter = str(row['lletra_inici']).strip()
            if letter.upper() != 'X':
                number += letter
        full_street = f"{full_street} {number}"
    
    parts = [full_street]
    if pd.notna(row.get(barri_col)):
        parts.append(str(row[barri_col]).strip())
    if pd.notna(row.get(districte_col)):
        parts.append(str(row[districte_col]).strip())
    parts.extend(['Barcelona', 'Spain'])
    
    return re.sub(r'\s+', ' ', ', '.join(parts))


def get_color_for_count(count):
    if count >= 5:
        return COUNT_COLORS['5+'], '5+'
    elif count == 4:
        return COUNT_COLORS['4'], '4'
    elif count == 3:
        return COUNT_COLORS['3'], '3'
    elif count == 2:
        return COUNT_COLORS['2'], '2'
    else:
        return COUNT_COLORS['1'], '1'


def extract_street_info(address):
    if not address:
        return "Unknown", None
    parts = address.split(',')
    if not parts:
        return address, None
    
    street_part = parts[0].strip()
    numbers = re.findall(r'\b\d+\b', street_part)
    street_number = int(numbers[0]) if numbers else None
    
    words = street_part.split()
    street_words = [w for w in words if not w.isdigit()]
    street_name = ' '.join(street_words) if street_words else street_part
    
    if street_name and street_name[0].islower():
        street_name = 'C. ' + street_name
    
    return street_name, street_number


def wrap_text(text, max_chars=65):
    words = text.split()
    lines = []
    current_line = ''
    for word in words:
        if len(current_line) + len(word) + 1 <= max_chars:
            current_line += (' ' if current_line else '') + word
        else:
            if current_line:
                lines.append(current_line)
            current_line = word
    if current_line:
        lines.append(current_line)
    return '\n'.join(lines)


# ============================================================================
# PHASE 1: DATA PROCESSING
# ============================================================================

def process_excel_data(df, street_col, num_col, barri_col, districte_col, progress_callback=None):
    if progress_callback:
        progress_callback(0.1, "Analyzing columns...")
    
    df_processed = df.copy()
    
    if progress_callback:
        progress_callback(0.3, "Creating complete addresses...")
    
    df_processed['complete_address'] = df_processed.apply(
        lambda row: create_complete_address(row, street_col, num_col, barri_col, districte_col),
        axis=1
    )
    
    if progress_callback:
        progress_callback(1.0, "Address processing complete")
    
    return df_processed


# ============================================================================
# PHASE 2: GEOCODING
# ============================================================================

def geocode_addresses(df, progress_callback=None):
    latitudes = []
    longitudes = []
    total = len(df)
    
    for idx, row in df.iterrows():
        address = row['complete_address']
        
        try:
            location = ox.geocode(address)
            latitudes.append(location[0])
            longitudes.append(location[1])
        except Exception:
            latitudes.append(None)
            longitudes.append(None)
        
        if progress_callback:
            progress = (idx + 1) / total
            progress_callback(progress, f"Geocoding {idx + 1}/{total}...")
        
        time.sleep(1.1)
    
    df['latitude'] = latitudes
    df['longitude'] = longitudes
    
    df_valid = df[df['latitude'].notna()].copy()
    
    gdf = gpd.GeoDataFrame(
        df_valid,
        geometry=[Point(lon, lat) for lon, lat in zip(df_valid['longitude'], df_valid['latitude'])],
        crs='EPSG:4326'
    )
    
    return gdf


# ============================================================================
# PHASE 3: CLUSTERING
# ============================================================================

def create_cluster_from_points(points):
    x_coords = [p['x'] for p in points]
    y_coords = [p['y'] for p in points]
    
    centroid_x = np.mean(x_coords)
    centroid_y = np.mean(y_coords)
    
    street_names = list(set(p['street_name'] for p in points))
    street_nums = [p['street_num'] for p in points if p['street_num'] is not None]
    realistic_nums = [num for num in street_nums if num < 1000]
    
    if len(street_names) == 1:
        street_name = street_names[0]
        if realistic_nums:
            unique_nums = set(realistic_nums)
            if len(unique_nums) == 1:
                street_display = f"{street_name}, {realistic_nums[0]}"
            else:
                min_num = min(realistic_nums)
                max_num = max(realistic_nums)
                if max_num - min_num <= 50:
                    street_display = f"{street_name}, {min_num}-{max_num}"
                else:
                    street_display = f"{street_name}, {min_num}"
        else:
            street_display = street_name
    else:
        street_display = ', '.join(sorted(street_names)[:3])
        if len(street_names) > 3:
            street_display += f" + {len(street_names)-3} mes"
    
    return {
        'centroid': (centroid_x, centroid_y),
        'count': len(points),
        'street': street_display,
        'points': points
    }


def merge_nearby_points(gdf_page, merge_radius=25):
    if gdf_page.empty:
        return []
    
    coords = np.array([[geom.x, geom.y] for geom in gdf_page.geometry])
    clustering = DBSCAN(eps=merge_radius, min_samples=1).fit(coords)
    
    spatial_clusters = defaultdict(list)
    for idx, cluster_label in enumerate(clustering.labels_):
        spatial_clusters[cluster_label].append(idx)
    
    final_clusters = []
    
    for cluster_indices in spatial_clusters.values():
        cluster_rows = gdf_page.iloc[cluster_indices]
        
        street_info = []
        for idx, row in zip(cluster_indices, cluster_rows.iterrows()):
            _, row_data = row
            street_name, street_num = extract_street_info(row_data.get('complete_address', ''))
            street_info.append({
                'idx': idx,
                'street_name': street_name,
                'street_num': street_num,
                'x': row_data.geometry.x,
                'y': row_data.geometry.y,
                'geometry': row_data.geometry
            })
        
        by_street = defaultdict(list)
        for info in street_info:
            by_street[info['street_name']].append(info)
        
        for street_name, points in by_street.items():
            points_with_nums = [p for p in points if p['street_num'] is not None]
            points_without_nums = [p for p in points if p['street_num'] is None]
            
            if points_with_nums:
                points_with_nums.sort(key=lambda p: p['street_num'])
                current_group = [points_with_nums[0]]
                
                for point in points_with_nums[1:]:
                    last_num = current_group[-1]['street_num']
                    curr_num = point['street_num']
                    
                    if abs(curr_num - last_num) <= 2:
                        current_group.append(point)
                    else:
                        final_clusters.append(create_cluster_from_points(current_group))
                        current_group = [point]
                
                if current_group:
                    final_clusters.append(create_cluster_from_points(current_group))
            
            for point in points_without_nums:
                final_clusters.append(create_cluster_from_points([point]))
    
    return final_clusters


def perform_zone_clustering(gdf, n_zones=6):
    gdf_web = gdf.to_crs(epsg=3857)
    coords = np.array([[geom.x, geom.y] for geom in gdf_web.geometry])
    
    n_clusters = min(n_zones, len(gdf_web))
    kmeans = KMeans(n_clusters=n_clusters, random_state=0, n_init=10)
    gdf_web['zone'] = kmeans.fit_predict(coords)
    
    return gdf_web


# ============================================================================
# OVERVIEW MAP FUNCTIONS (with MTM basemap and background layers)
# ============================================================================

def crop_white_borders(img, threshold=250):
    """Crop white borders from an image."""
    if len(img.shape) == 3:
        if img.shape[2] == 4:
            gray = np.mean(img[:, :, :3], axis=2)
        else:
            gray = np.mean(img, axis=2)
    else:
        gray = img
    
    if gray.max() <= 1.0:
        threshold = threshold / 255.0
    
    non_white_rows = np.where(np.min(gray, axis=1) < threshold)[0]
    non_white_cols = np.where(np.min(gray, axis=0) < threshold)[0]
    
    if len(non_white_rows) == 0 or len(non_white_cols) == 0:
        return img
    
    top = non_white_rows[0]
    bottom = non_white_rows[-1] + 1
    left = non_white_cols[0]
    right = non_white_cols[-1] + 1
    
    return img[top:bottom, left:right]


def generate_overview_reference_image(gdf_full, gdf_zone, zone_id, carril_bici_gdf=None, carrers_30_gdf=None):
    """
    Generate a reference overview map image showing zone location within Barcelona.
    Uses MTM basemap and includes carril_bici and carrers_30 layers like main map.
    """
    try:
        gdf_full_web = gdf_full.to_crs(epsg=3857)
        gdf_zone_web = gdf_zone.to_crs(epsg=3857)
        
        # Create figure for overview
        fig_overview = plt.figure(figsize=(6, 6))
        ax_overview = fig_overview.add_subplot(111)
        
        # Get bounds
        full_bounds = gdf_full_web.total_bounds
        zone_bounds = gdf_zone_web.total_bounds
        
        # Add padding
        pad_x = (full_bounds[2] - full_bounds[0]) * 0.08
        pad_y = (full_bounds[3] - full_bounds[1]) * 0.08
        
        ax_overview.set_xlim(full_bounds[0] - pad_x, full_bounds[2] + pad_x)
        ax_overview.set_ylim(full_bounds[1] - pad_y, full_bounds[3] + pad_y)
        ax_overview.set_aspect('equal', adjustable='box')
        
        # Add MTM basemap (same as main map)
        try:
            cx.add_basemap(
                ax_overview,
                source=BARCELONA_MTM_TILES,
                zoom='auto',
                alpha=0.8
            )
        except Exception:
            try:
                cx.add_basemap(
                    ax_overview,
                    source=cx.providers.CartoDB.Positron,
                    zoom='auto',
                    alpha=0.7
                )
            except Exception:
                ax_overview.set_facecolor('#f5f5f5')
        
        # Add background layers (same style as main map)
        if carrers_30_gdf is not None:
            try:
                carrers_30_web = carrers_30_gdf.to_crs(epsg=3857)
                carrers_30_web.plot(ax=ax_overview, color='#66B3FF', linewidth=0.8, alpha=0.5, zorder=2)
            except Exception:
                pass
        
        if carril_bici_gdf is not None:
            try:
                carril_bici_web = carril_bici_gdf.to_crs(epsg=3857)
                carril_bici_web.plot(ax=ax_overview, color='#FF99AA', linewidth=1.0, alpha=0.4, zorder=3)
            except Exception:
                pass
        
        # Plot all points in light gray
        gdf_full_web.plot(ax=ax_overview, color='#999999', markersize=6, alpha=0.3, zorder=4)
        
        # Plot current zone points in red
        gdf_zone_web.plot(ax=ax_overview, color='#d62728', markersize=12, alpha=0.8, zorder=5)
        
        # Draw rectangle around current zone
        rect = Rectangle(
            (zone_bounds[0], zone_bounds[1]),
            zone_bounds[2] - zone_bounds[0],
            zone_bounds[3] - zone_bounds[1],
            fill=False,
            edgecolor='#d62728',
            linewidth=2.5,
            zorder=6
        )
        ax_overview.add_patch(rect)
        
        # Add zone label
        ax_overview.text(
            0.5, 0.97,
            f'Zona {zone_id + 1}',
            transform=ax_overview.transAxes,
            ha='center', va='top',
            fontsize=10,
            fontweight='bold',
            color='#333333',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#cccccc', alpha=0.9)
        )
        
        ax_overview.set_xticks([])
        ax_overview.set_yticks([])
        ax_overview.axis('off')
        
        plt.tight_layout(pad=0)
        
        # Convert figure to image array
        fig_overview.canvas.draw()
        img = np.frombuffer(fig_overview.canvas.tostring_rgb(), dtype=np.uint8)
        img = img.reshape(fig_overview.canvas.get_width_height()[::-1] + (3,))
        
        plt.close(fig_overview)
        
        # Crop white borders
        img = crop_white_borders(img)
        
        return img
        
    except Exception as e:
        print(f"Error generating overview: {e}")
        return None


def add_overview_map_to_panel(ax_box, ax_map, fig, zone_id, gdf_zone, gdf_full, 
                               carril_bici_gdf=None, carrers_30_gdf=None, overview_height=0.45):
    """
    Add an overview reference map to the right panel, bottom-aligned with the main map.
    EXACT positioning logic from notebook.
    """
    try:
        # Generate the overview image with MTM and background layers
        img = generate_overview_reference_image(gdf_full, gdf_zone, zone_id, carril_bici_gdf, carrers_30_gdf)
        
        if img is None:
            return
        
        # Get image dimensions and aspect ratio
        img_height_px, img_width_px = img.shape[:2]
        img_aspect = img_width_px / img_height_px
        
        # Get the bounding boxes in figure coordinates
        fig.canvas.draw()
        
        # Get main map position
        map_bbox = ax_map.get_position()
        map_bottom = map_bbox.y0
        
        # Get right panel position
        box_bbox = ax_box.get_position()
        box_left = box_bbox.x0
        box_width = box_bbox.width
        box_height = box_bbox.height
        
        # Calculate display dimensions (EXACT from notebook)
        x_margin_fig = box_width * 0.005
        display_width_fig = box_width - 2 * x_margin_fig
        
        # Calculate height maintaining aspect ratio
        fig_width, fig_height = fig.get_size_inches()
        fig_aspect = fig_width / fig_height
        
        display_height_fig = (display_width_fig / img_aspect) * fig_aspect
        
        # Check if height exceeds available space
        max_height_fig = box_height * overview_height
        if display_height_fig > max_height_fig:
            display_height_fig = max_height_fig
            display_width_fig = (display_height_fig * img_aspect) / fig_aspect
        
        # Position: bottom aligned with main map, centered horizontally
        x_pos = box_left + (box_width - display_width_fig) / 2
        y_pos = map_bottom
        
        # Create axes at exact figure coordinates
        ax_inset = fig.add_axes([x_pos, y_pos, display_width_fig, display_height_fig])
        
        # Display the image
        ax_inset.imshow(img, aspect='equal')
        
        # Remove all axes elements
        ax_inset.set_xticks([])
        ax_inset.set_yticks([])
        ax_inset.axis('off')
        
    except Exception as e:
        print(f"Warning: Could not add overview map: {e}")


# ============================================================================
# PHASE 4: MAP GENERATION (exact from notebook)
# ============================================================================

def generate_zone_map(gdf_zone, zone_id, gdf_full, merge_radius=25, global_max_count=None,
                      carril_bici_gdf=None, carrers_30_gdf=None):
    """Generate a single zone map - EXACT from notebook."""
    
    gdf_zone_web = gdf_zone.to_crs(epsg=3857)
    
    clusters = merge_nearby_points(gdf_zone_web, merge_radius)
    
    if not clusters:
        return None
    
    # Assign colors
    for cluster in clusters:
        cluster['color'], cluster['color_label'] = get_color_for_count(cluster['count'])
    
    # Sort by category priority
    def sort_key(cluster):
        return COUNT_ORDER.index(cluster['color_label'])
    
    clusters.sort(key=sort_key)
    
    for idx, cluster in enumerate(clusters):
        cluster['number'] = idx + 1
    
    if global_max_count is None:
        global_max_count = max(c['count'] for c in clusters)
    
    # Get page bounds
    page_bounds = gdf_zone_web.total_bounds
    minx, miny, maxx, maxy = page_bounds
    
    x_range = maxx - minx
    y_range = maxy - miny
    
    if x_range == 0:
        x_range = 1
        minx -= 0.5
        maxx += 0.5
    if y_range == 0:
        y_range = 1
        miny -= 0.5
        maxy += 0.5
    
    # Calculate proper aspect ratio (EXACT from notebook)
    map_ar = MAP_FRAC * A3_WIDTH_IN / A3_HEIGHT_IN
    data_ar = x_range / y_range
    cx_center = (minx + maxx) / 2
    cy_center = (miny + maxy) / 2
    
    if data_ar < map_ar:
        new_x_range = y_range * map_ar
        minx = cx_center - new_x_range / 2
        maxx = cx_center + new_x_range / 2
    else:
        new_y_range = x_range / map_ar
        miny = cy_center - new_y_range / 2
        maxy = cy_center + new_y_range / 2
    
    pad_x = (maxx - minx) * 0.05
    pad_y = (maxy - miny) * 0.05
    minx -= pad_x
    maxx += pad_x
    miny -= pad_y
    maxy += pad_y
    
    # Create figure (EXACT from notebook)
    fig = plt.figure(figsize=(A3_WIDTH_IN, A3_HEIGHT_IN))
    gs = gridspec.GridSpec(1, 2, width_ratios=[MAP_FRAC, BOX_FRAC], wspace=0.01)
    ax_map = fig.add_subplot(gs[0])
    ax_map.set_xlim(minx, maxx)
    ax_map.set_ylim(miny, maxy)
    ax_map.set_aspect('equal', adjustable='box')
    
    # Add basemap
    try:
        cx.add_basemap(ax_map, source=BARCELONA_MTM_TILES, zoom='auto', alpha=0.8, zorder=1)
    except Exception:
        try:
            cx.add_basemap(ax_map, source=cx.providers.CartoDB.Positron, zoom='auto', alpha=0.8, zorder=1)
        except Exception:
            pass
    
    # Add background layers (EXACT from notebook)
    if carrers_30_gdf is not None:
        try:
            carrers_30_web = carrers_30_gdf.to_crs(epsg=3857)
            carrers_30_web.plot(ax=ax_map, color='#66B3FF', linewidth=1.5, alpha=0.6, zorder=2)
        except Exception:
            pass
    
    if carril_bici_gdf is not None:
        try:
            carril_bici_web = carril_bici_gdf.to_crs(epsg=3857)
            carril_bici_web.plot(ax=ax_map, color='#FF99AA', linewidth=1.8, alpha=0.4, zorder=3)
        except Exception:
            pass
    
    # Plot clusters (EXACT from notebook)
    min_markersize = 65
    max_markersize = 600
    min_fontsize = 4
    max_fontsize = 4
    
    for cluster in clusters:
        centroid_x, centroid_y = cluster['centroid']
        count_ratio = cluster['count'] / global_max_count
        markersize = min_markersize + (max_markersize - min_markersize) * count_ratio
        
        ax_map.scatter(
            centroid_x, centroid_y,
            s=markersize,
            color=cluster['color'],
            marker='o',
            alpha=0.8,
            edgecolors='white',
            linewidths=1.5,
            zorder=5
        )
        
        fontsize = min_fontsize + (max_fontsize - min_fontsize) * count_ratio
        ax_map.text(
            centroid_x, centroid_y,
            str(cluster['number']),
            ha='center', va='center',
            fontsize=fontsize,
            fontweight='bold',
            color='white',
            zorder=6
        )
    
    # Title
    ax_map.text(
        0.02, 0.98,
        f'Concentracio IRIS - Zona {zone_id + 1}',
        transform=ax_map.transAxes,
        fontsize=8,
        fontweight='bold',
        verticalalignment='top',
        horizontalalignment='left',
        bbox=dict(
            boxstyle='square,pad=0.5',
            facecolor='white',
            edgecolor='black',
            linewidth=1,
            alpha=0.95
        ),
        zorder=10
    )
    
    # Legend
    present_categories = set(c['color_label'] for c in clusters)
    legend_items = []
    for cat in COUNT_ORDER:
        if cat in present_categories:
            legend_items.append(
                Patch(facecolor=COUNT_COLORS[cat], edgecolor='white', linewidth=1,
                      label=f'{cat} IRIS' if cat != '5+' else '5+ IRIS')
            )
    
    if legend_items:
        ax_map.legend(
            handles=legend_items,
            loc='lower right',
            fontsize=5,
            title_fontsize=6,
            framealpha=0.95,
            fancybox=True,
            shadow=True,
            bbox_to_anchor=(0.98, 0.02)
        )
    
    ax_map.set_xticks([])
    ax_map.set_yticks([])
    ax_map.set_axis_off()
    
    # RIGHT PANEL
    ax_box = fig.add_subplot(gs[1])
    ax_box.set_xlim(0, 1)
    ax_box.set_ylim(0, 1)
    ax_box.set_aspect('auto')
    ax_box.set_facecolor('white')
    ax_box.add_patch(Rectangle((0, 0), 1, 1, fill=False, edgecolor='white', linewidth=2))
    
    clusters_by_category = defaultdict(list)
    for cluster in clusters:
        clusters_by_category[cluster['color_label']].append(cluster)
    
    num_clusters = len(clusters)
    
    # Legend sizing (EXACT from notebook)
    overview_map_height = 0.45
    base_line_height = 0.015
    header_extra = 0.006
    category_spacing = 0.01
    text_fontsize = 5.0
    header_fontsize = 5.5
    
    if num_clusters > 20:
        base_line_height = 0.015
        category_spacing = 0.008
        text_fontsize = 5.0
    
    y_start = 0.98
    y_bottom_limit = 0.02
    current_y = y_start
    
    first_category = True
    for cat in COUNT_ORDER:
        if cat not in clusters_by_category:
            continue
        
        cat_clusters = clusters_by_category[cat]
        cat_color = COUNT_COLORS[cat]
        cat_total = sum(c['count'] for c in cat_clusters)
        
        if not first_category:
            current_y -= category_spacing
        first_category = False
        
        if current_y - (base_line_height + header_extra) < y_bottom_limit:
            break
        
        ax_box.add_patch(Rectangle(
            (0.02, current_y - 0.012),
            0.96, 0.024,
            facecolor='white',
            edgecolor='none',
            alpha=1.0,
            zorder=10
        ))
        ax_box.add_patch(Rectangle(
            (0.05, current_y - 0.010),
            0.90, 0.020,
            facecolor=cat_color,
            edgecolor='none',
            alpha=0.2,
            zorder=11
        ))
        
        header_label = f"5+ IRIS ({cat_total} total)" if cat == '5+' else f"{cat} IRIS ({cat_total} total)"
        
        ax_box.text(
            0.18, current_y,
            header_label,
            transform=ax_box.transAxes,
            ha='left', va='center',
            fontsize=header_fontsize,
            fontweight='bold',
            color=cat_color,
            zorder=12
        )
        current_y -= (base_line_height + header_extra)
        
        for cluster in cat_clusters:
            num_lines = wrap_text(cluster['street'], max_chars=60).count('\n') + 1
            entry_height = base_line_height * max(1, num_lines * 0.85)
            
            if current_y - entry_height < y_bottom_limit:
                break
            
            cluster_num = cluster['number']
            
            ax_box.add_patch(Rectangle(
                (0.02, current_y - entry_height/2 - 0.002),
                0.96, entry_height + 0.004,
                facecolor='white',
                edgecolor='none',
                alpha=1.0,
                zorder=10
            ))
            
            ax_box.text(
                0.10, current_y,
                str(cluster_num),
                transform=ax_box.transAxes,
                ha='center', va='center',
                fontsize=5,
                fontweight='bold',
                color='black',
                zorder=12
            )
            
            street_text = cluster['street']
            wrapped_text = wrap_text(street_text, max_chars=60)
            
            ax_box.text(
                0.18, current_y,
                wrapped_text,
                transform=ax_box.transAxes,
                ha='left', va='center',
                fontsize=text_fontsize,
                zorder=12
            )
            
            current_y -= entry_height
    
    ax_box.set_xticks([])
    ax_box.set_yticks([])
    ax_box.set_axis_off()
    
    # Apply tight_layout BEFORE adding overview map
    plt.tight_layout(pad=0)
    
    # ADD OVERVIEW MAP with MTM and background layers
    add_overview_map_to_panel(
        ax_box, ax_map, fig,
        zone_id, gdf_zone, gdf_full,
        carril_bici_gdf=carril_bici_gdf,
        carrers_30_gdf=carrers_30_gdf,
        overview_height=overview_map_height
    )
    
    return fig, clusters


def generate_all_maps(gdf, n_zones=6, merge_radius=25, carril_bici_gdf=None, carrers_30_gdf=None, progress_callback=None):
    """Generate maps for all zones."""
    gdf_zoned = perform_zone_clustering(gdf, n_zones)
    
    # First pass: calculate global max count
    all_clusters_initial = []
    for zone_id in range(n_zones):
        gdf_zone = gdf_zoned[gdf_zoned['zone'] == zone_id]
        if len(gdf_zone) > 0:
            gdf_zone_web = gdf_zone.to_crs(epsg=3857)
            page_clusters = merge_nearby_points(gdf_zone_web, merge_radius=merge_radius)
            all_clusters_initial.extend(page_clusters)
    
    global_max_count = max(c['count'] for c in all_clusters_initial) if all_clusters_initial else 1
    
    figures = []
    all_clusters = []
    
    for zone_id in range(n_zones):
        if progress_callback:
            progress_callback(zone_id / n_zones, f"Generating map for Zone {zone_id + 1}...")
        
        gdf_zone = gdf_zoned[gdf_zoned['zone'] == zone_id]
        
        if len(gdf_zone) == 0:
            continue
        
        result = generate_zone_map(
            gdf_zone, 
            zone_id, 
            gdf_zoned,
            merge_radius,
            global_max_count=global_max_count,
            carril_bici_gdf=carril_bici_gdf,
            carrers_30_gdf=carrers_30_gdf
        )
        
        if result is not None:
            fig, clusters = result
            figures.append((zone_id, fig))
            all_clusters.extend(clusters)
    
    if progress_callback:
        progress_callback(1.0, "All maps generated")
    
    return figures, all_clusters


# ============================================================================
# PHASE 5: POWERPOINT GENERATION (FIXED - minimal white space)
# ============================================================================

def create_powerpoint(figures, output_path, dpi=300, carril_bici_gdf=None, carrers_30_gdf=None,
                      gdf=None, n_zones=6, merge_radius=25, progress_callback=None):
    """Create PowerPoint from generated figures - image fills entire slide."""
    
    # Regenerate figures fresh for PowerPoint
    if gdf is not None:
        figures_fresh, _ = generate_all_maps(
            gdf, n_zones=n_zones, merge_radius=merge_radius,
            carril_bici_gdf=carril_bici_gdf, carrers_30_gdf=carrers_30_gdf
        )
    else:
        figures_fresh = figures
    
    prs = Presentation()
    # Use widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    temp_images = []
    
    for idx, (zone_id, fig) in enumerate(figures_fresh):
        if progress_callback:
            progress_callback(idx / len(figures_fresh), f"Adding slide {idx + 1}...")
        
        # Save figure to temp file
        temp_path = tempfile.mktemp(suffix='.png')
        
        # Save with exact figure size, no extra padding
        fig.savefig(
            temp_path, 
            dpi=dpi, 
            facecolor='white',
            edgecolor='none',
            bbox_inches='tight',
            pad_inches=0
        )
        plt.close(fig)
        temp_images.append(temp_path)
        
        # Add slide
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Load image to get dimensions
        with Image.open(temp_path) as img:
            img_width, img_height = img.size
        
        # Calculate scaling to FILL slide (may crop edges slightly)
        img_aspect = img_width / img_height
        slide_aspect = prs.slide_width / prs.slide_height
        
        if img_aspect > slide_aspect:
            # Image is wider - scale to height, center horizontally
            pic_height = prs.slide_height
            pic_width = pic_height * img_aspect
            left = (prs.slide_width - pic_width) / 2
            top = Inches(0)
        else:
            # Image is taller - scale to width, center vertically
            pic_width = prs.slide_width
            pic_height = pic_width / img_aspect
            left = Inches(0)
            top = (prs.slide_height - pic_height) / 2
        
        slide.shapes.add_picture(
            temp_path,
            left, top,
            width=pic_width,
            height=pic_height
        )
    
    prs.save(output_path)
    
    # Cleanup
    for temp_path in temp_images:
        try:
            os.remove(temp_path)
        except Exception:
            pass
    
    if progress_callback:
        progress_callback(1.0, "PowerPoint created")
    
    return output_path


# ============================================================================
# STREAMLIT UI
# ============================================================================

def main():
    st.markdown("# IRIS Pipeline")
    st.markdown("Barcelona Bicycle Incidents Analysis | Excel to Maps to PowerPoint")
    
    # Initialize session state
    if 'df_processed' not in st.session_state:
        st.session_state.df_processed = None
    if 'gdf' not in st.session_state:
        st.session_state.gdf = None
    if 'figures' not in st.session_state:
        st.session_state.figures = None
    if 'pptx_data' not in st.session_state:
        st.session_state.pptx_data = None
    if 'map_images' not in st.session_state:
        st.session_state.map_images = None
    if 'carril_bici_gdf' not in st.session_state:
        st.session_state.carril_bici_gdf = None
    if 'carrers_30_gdf' not in st.session_state:
        st.session_state.carrers_30_gdf = None
    
    # Sidebar
    with st.sidebar:
        st.markdown("### Configuration")
        st.markdown("---")
        
        st.markdown("#### Clustering Parameters")
        merge_radius = st.slider("Merge radius (meters)", 10, 100, 25, 5)
        n_zones = st.slider("Number of zones", 2, 10, 6)
        
        st.markdown("---")
        
        st.markdown("#### Output Settings")
        dpi = st.selectbox("Map DPI", [150, 200, 300], index=2)
        
        st.markdown("---")
        
        st.markdown("#### Background Layers")
        carril_file = st.file_uploader("CARRIL_BICI.geojson", type=['geojson', 'json'], key="carril")
        carrers_file = st.file_uploader("CARRERS_30.geojson", type=['geojson', 'json'], key="carrers")
        
        if carril_file is not None:
            st.session_state.carril_bici_gdf = gpd.read_file(carril_file)
            st.success("Carril bici loaded")
        
        if carrers_file is not None:
            st.session_state.carrers_30_gdf = gpd.read_file(carrers_file)
            st.success("Carrers 30 loaded")
        
        st.markdown("---")
        
        st.markdown("### Pipeline Status")
        phases = [
            ("1. Data Upload", st.session_state.df_processed is not None),
            ("2. Geocoding", st.session_state.gdf is not None),
            ("3. Map Generation", st.session_state.figures is not None),
            ("4. PowerPoint Export", st.session_state.pptx_data is not None),
        ]
        
        for phase_name, completed in phases:
            status = "[Done]" if completed else "[Pending]"
            st.markdown(f"{status} {phase_name}")
    
    # Main content
    tab1, tab2, tab3, tab4 = st.tabs([
        "1. Upload Data",
        "2. Geocode",
        "3. Generate Maps",
        "4. Export"
    ])
    
    # TAB 1
    with tab1:
        st.markdown("### Upload Excel File")
        
        uploaded_file = st.file_uploader(
            "Choose Excel file",
            type=['xlsx', 'xls'],
            help="Excel file with incident data"
        )
        
        if uploaded_file is not None:
            df = pd.read_excel(uploaded_file)
            
            st.markdown("### Data Preview")
            st.markdown(df.head(10).to_html(index=False), unsafe_allow_html=True)
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Rows", len(df))
            with col2:
                st.metric("Columns", len(df.columns))
            with col3:
                st.metric("Unique Streets", df['carrer'].nunique() if 'carrer' in df.columns else "N/A")
            with col4:
                st.metric("Districts", df['districte'].nunique() if 'districte' in df.columns else "N/A")
            
            st.markdown("### Column Mapping")
            
            col1, col2 = st.columns(2)
            with col1:
                street_col = st.selectbox("Street name column", df.columns.tolist(), 
                                         index=df.columns.tolist().index('carrer') if 'carrer' in df.columns else 0)
                barri_col = st.selectbox("Neighborhood column", df.columns.tolist(),
                                        index=df.columns.tolist().index('barri') if 'barri' in df.columns else 0)
            
            with col2:
                num_col = st.selectbox("Street number column", df.columns.tolist(),
                                      index=df.columns.tolist().index('numero_inici') if 'numero_inici' in df.columns else 0)
                districte_col = st.selectbox("District column", df.columns.tolist(),
                                            index=df.columns.tolist().index('districte') if 'districte' in df.columns else 0)
            
            if st.button("Process Addresses", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                def update_progress(progress, message):
                    progress_bar.progress(progress)
                    status_text.text(message)
                
                st.session_state.df_processed = process_excel_data(
                    df, street_col, num_col, barri_col, districte_col, update_progress
                )
                
                st.success(f"Processed {len(st.session_state.df_processed)} addresses successfully.")
    
    # TAB 2
    with tab2:
        st.markdown("### Geocode Addresses")
        
        if st.session_state.df_processed is None:
            st.warning("Please upload and process data first (Tab 1)")
        else:
            st.markdown(f"Ready to geocode **{len(st.session_state.df_processed)}** addresses.")
            
            st.markdown("---")
            st.markdown("#### Or Upload Pre-Geocoded GeoJSON")
            
            geojson_file = st.file_uploader(
                "Upload existing GeoJSON (skip geocoding)",
                type=['geojson', 'json'],
                key="geojson_upload"
            )
            
            if geojson_file is not None:
                gdf = gpd.read_file(geojson_file)
                st.session_state.gdf = gdf.to_crs(epsg=4326)
                st.success(f"Loaded {len(st.session_state.gdf)} features from GeoJSON.")
            
            st.markdown("---")
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("Start Geocoding", use_container_width=True, disabled=st.session_state.gdf is not None):
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    def update_progress(progress, message):
                        progress_bar.progress(progress)
                        status_text.text(message)
                    
                    with st.spinner("Geocoding in progress..."):
                        st.session_state.gdf = geocode_addresses(
                            st.session_state.df_processed.copy(),
                            update_progress
                        )
                    
                    st.success(f"Geocoded {len(st.session_state.gdf)} addresses")
            
            with col2:
                if st.session_state.gdf is not None:
                    geojson_str = st.session_state.gdf.to_json()
                    st.download_button(
                        "Download GeoJSON",
                        data=geojson_str,
                        file_name="incidents_geocoded.geojson",
                        mime="application/json",
                        use_container_width=True
                    )
    
    # TAB 3
    with tab3:
        st.markdown("### Generate Zone Maps")
        
        if st.session_state.gdf is None:
            st.warning("Please complete geocoding first (Tab 2)")
        else:
            bg_status = []
            if st.session_state.carril_bici_gdf is not None:
                bg_status.append("Carril Bici")
            if st.session_state.carrers_30_gdf is not None:
                bg_status.append("Carrers 30")
            
            bg_text = f" with {', '.join(bg_status)}" if bg_status else " (no background layers)"
            
            st.markdown(f"Generate maps divided into **{n_zones} zones** with merge radius **{merge_radius}m**{bg_text}")
            
            if st.button("Generate Maps", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                def update_progress(progress, message):
                    progress_bar.progress(progress)
                    status_text.text(message)
                
                with st.spinner("Generating maps..."):
                    st.session_state.figures, clusters = generate_all_maps(
                        st.session_state.gdf,
                        n_zones=n_zones,
                        merge_radius=merge_radius,
                        carril_bici_gdf=st.session_state.carril_bici_gdf,
                        carrers_30_gdf=st.session_state.carrers_30_gdf,
                        progress_callback=update_progress
                    )
                    
                    # Save as images for preview
                    st.session_state.map_images = []
                    for zone_id, fig in st.session_state.figures:
                        buf = BytesIO()
                        fig.savefig(buf, format='png', dpi=100, facecolor='white', bbox_inches='tight', pad_inches=0)
                        buf.seek(0)
                        st.session_state.map_images.append((zone_id, buf.getvalue()))
                
                st.success(f"Generated {len(st.session_state.figures)} zone maps.")
            
            if st.session_state.map_images is not None:
                st.markdown("### Map Preview")
                
                for zone_id, img_data in st.session_state.map_images:
                    with st.expander(f"Zone {zone_id + 1}", expanded=(zone_id == 0)):
                        st.image(img_data, use_container_width=True)
    
    # TAB 4
    with tab4:
        st.markdown("### Export to PowerPoint")
        
        if st.session_state.figures is None:
            st.warning("Please generate maps first (Tab 3)")
        else:
            st.markdown(f"Create PowerPoint with **{len(st.session_state.figures)} slides** at **{dpi} DPI**")
            
            if st.button("Create PowerPoint", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                def update_progress(progress, message):
                    progress_bar.progress(progress)
                    status_text.text(message)
                
                with st.spinner("Creating PowerPoint..."):
                    temp_pptx = tempfile.mktemp(suffix='.pptx')
                    create_powerpoint(
                        st.session_state.figures, 
                        temp_pptx, 
                        dpi=dpi,
                        carril_bici_gdf=st.session_state.carril_bici_gdf,
                        carrers_30_gdf=st.session_state.carrers_30_gdf,
                        gdf=st.session_state.gdf,
                        n_zones=n_zones,
                        merge_radius=merge_radius,
                        progress_callback=update_progress
                    )
                    
                    with open(temp_pptx, 'rb') as f:
                        st.session_state.pptx_data = f.read()
                    
                    os.remove(temp_pptx)
                
                st.success("PowerPoint created successfully.")
            
            if st.session_state.pptx_data is not None:
                st.markdown("---")
                
                timestamp = datetime.now().strftime("%Y%m%d_%H%M")
                filename = f"iris_incidents_{timestamp}.pptx"
                
                st.download_button(
                    "Download PowerPoint",
                    data=st.session_state.pptx_data,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                st.markdown("---")
                st.markdown("**Pipeline Complete.** Download your PowerPoint above.")
    
    # Footer
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #888; font-size: 0.75rem;'>"
        "IRIS Pipeline | Barcelona Mobility Analysis"
        "</div>",
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    main()
